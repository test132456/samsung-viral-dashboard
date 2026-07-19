"""블로그 작성 가이드(.pptx) 파서 & 원고-가이드 준수 체크. 순수 함수.

extract_text(pptx_bytes)  가이드 PPT 전체 텍스트(표는 'a | b | c'로)
parse_guide(text)         가이드에서 금지표현·필수 해시태그 추출
check(manuscript, guide)  원고를 가이드 기준으로 점검
"""
from __future__ import annotations
import io
import re

_QUOTES = "‘’\"'＇「」“”"
_HEADER_TERMS = {"불가 표현", "불가표현", "불가 표현 ", "예시", "사유"}


def extract_text(pptx_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    lines = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                lines.append(sh.text_frame.text.strip())
            if sh.has_table:
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
    return "\n".join(lines)


_RIDER_SECTION_RE = re.compile(r"\[[^\]]*담보명[^\]]*\]")
# 담보명임을 판별하는 핵심 토큰 (법조문 조각 '이 특약/제1조(특약' 등 걸러내기)
_RIDER_TOKENS = ("항공기", "수하물", "휴대품", "식중독", "여권", "질병", "상해",
                 "실손", "배상", "치료", "지연", "결항", "분실", "손해", "재발급", "보상금")
# 가이드 '메인 담보명' 정식 특약명 (가이드 PPT 미인식 시 기본값)
DEFAULT_RIDERS = [
    "항공기 지연 결항 보상(지수형)(국내 출국) 특약",
    "항공기 지연 결항 보상(지수형)(국내 출국 제외) 특약",
    "수하물 지연(6시간 이상)·손실 추가비용 특약",
    "여행중 휴대품 손해(분실제외) 특약",
    "여행중 식중독 보상금(2일이상 입원)특약",
    "여행중 여권분실 재발급비용 특약",
]


_RIDER_SEP = "[" + chr(10) + chr(13) + chr(11) + chr(12) + "]+"  # 개행·수직탭 등
_RIDER_NOISE = ("확인", "정확", "소구", "담보명", "약관에서", "포인트", "메인")


def _extract_riders(text: str) -> list[str]:
    """가이드 '메인 담보명' 표에서 정식 특약명 추출.
    표 각 행의 첫 칸(| 기준)이 '…특약'으로 끝나는 것만 담보명으로 본다.
    → '출국 항공기 지연 손해 특약 제외'(제외로 끝남)·안내문·소구문구 등은 자동 배제.
    담보명 섹션 헤더/표가 없으면 [] (호출부에서 기본값으로 대체)."""
    if "담보명" not in (text or ""):
        return []
    out, seen = [], set()
    for raw in re.split(_RIDER_SEP, text):  # PPT 셀 줄바꿈=수직탭
        cell = raw.split("|")[0].strip().strip("★").strip()
        cell = re.sub(r"^\s*(?:제?\s*\d+\s*[.)관조항호]|[·•\-*])\s*", "", cell).strip()
        if not cell.endswith("특약"):           # 이름은 반드시 '특약'으로 끝남
            continue
        if not (4 <= len(cell) <= 45):
            continue
        if any(w in cell for w in _RIDER_NOISE):
            continue
        if not any(t in cell for t in _RIDER_TOKENS):
            continue
        if cell not in seen:
            seen.add(cell)
            out.append(cell)
    return out


def parse_guide(text: str) -> dict:
    text = text or ""
    # 필수 해시태그
    hashtags = []
    for h in re.findall(r"#[가-힣A-Za-z0-9]+", text):
        if h not in hashtags:
            hashtags.append(h)
    # 금지표현: '표현 불가' 섹션의 표 첫 칸
    banned = []
    if "표현 불가" in text:
        section = text.split("표현 불가", 1)[1]
        for line in section.splitlines():
            if "|" not in line:
                continue
            raw = line.split("|")[0].strip().strip(_QUOTES).strip()
            if not (2 <= len(raw) <= 25) or raw in _HEADER_TERMS:
                continue
            # '빈번하게, 자주'처럼 각 부분이 모두 2자 이상이면 개별 단어로 분리
            parts = [p.strip() for p in raw.split(",")]
            cands = parts if len(parts) > 1 and all(len(p) >= 2 for p in parts) else [raw]
            for c in cands:
                c = c.strip().strip("~").strip()
                if 2 <= len(c) <= 25 and c not in _HEADER_TERMS and c not in banned:
                    banned.append(c)
    return {"hashtags": hashtags, "banned": banned, "riders": _extract_riders(text)}


def check(manuscript: str, guide: dict) -> dict:
    ms = manuscript or ""
    banned_hits = [b for b in guide.get("banned", []) if b in ms]
    tags = guide.get("hashtags", [])
    tags_missing = [t for t in tags if t not in ms]
    tags_included = [t for t in tags if t in ms]
    kw_count = ms.count("해외여행보험")
    return {
        "banned_hits": banned_hits,
        "tags_total": len(tags), "tags_included": tags_included, "tags_missing": tags_missing,
        "keyword_count": kw_count,          # '해외여행보험' 등장 횟수 (가이드: 3~5개)
        "keyword_ok": 3 <= kw_count <= 5,
    }
