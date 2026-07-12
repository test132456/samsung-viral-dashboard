import io
from pptx import Presentation
from pptx.util import Inches
from core import guide


def _make_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9), Inches(1)).text_frame
    tb.text = "하단 해시태그: #삼성화재다이렉트 #해외여행보험 #항공지연"
    tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(9), Inches(0.5)).text_frame
    tb2.text = "[표현 불가 문구]"
    tbl = slide.shapes.add_table(3, 3, Inches(0.3), Inches(2.2), Inches(9), Inches(2)).table
    tbl.cell(0, 0).text, tbl.cell(0, 1).text, tbl.cell(0, 2).text = "불가 표현", "사유", "예시"
    tbl.cell(1, 0).text, tbl.cell(1, 1).text, tbl.cell(1, 2).text = "무조건", "삭제 필요", "예시문"
    tbl.cell(2, 0).text, tbl.cell(2, 1).text, tbl.cell(2, 2).text = "빈번하게, 자주", "근거 없음", "예시문"
    b = io.BytesIO()
    prs.save(b)
    return b.getvalue()


def test_parse_guide():
    g = guide.parse_guide(guide.extract_text(_make_pptx()))
    assert "#해외여행보험" in g["hashtags"] and "#삼성화재다이렉트" in g["hashtags"]
    assert "무조건" in g["banned"]
    assert "빈번하게" in g["banned"] and "자주" in g["banned"]  # 콤마 분리


def test_check_against_manuscript():
    g = guide.parse_guide(guide.extract_text(_make_pptx()))
    ms = "이 보험 무조건 좋아요. #삼성화재다이렉트 해외여행보험 해외여행보험 해외여행보험 안내."
    res = guide.check(ms, g)
    assert "무조건" in res["banned_hits"]
    assert "#해외여행보험" in res["tags_missing"]     # 원고에 이 태그 없음
    assert "#삼성화재다이렉트" in res["tags_included"]
    assert res["keyword_count"] == 3 and res["keyword_ok"] is True


def test_extract_riders_from_dambo_section():
    from core import guide
    text = ("[메인 담보명 및 소구 포인트] ★약관에서 확인되는 정확한 특약명\n"
            "특약명\n"
            "항공기 지연 결항 보상(지수형)(국내 출국) 특약\n"
            "수하물 지연(6시간 이상)·손실 추가비용 특약\n"
            "[표현 불가 문구]\n"
            "출국 항공기 지연 손해 특약 제외 | 보장 아님 |\n")
    g = guide.parse_guide(text)
    assert "항공기 지연 결항 보상(지수형)(국내 출국) 특약" in g["riders"]
    assert "수하물 지연(6시간 이상)·손실 추가비용 특약" in g["riders"]
    # 헤더/안내문/표현불가표의 특약은 제외
    assert not any("확인" in r for r in g["riders"])
    assert "출국 항공기 지연 손해 특약 제외" not in g["riders"]
    assert "특약명" not in g["riders"]


def test_extract_riders_none_without_section():
    from core import guide
    assert guide.parse_guide("담보명 헤더 없는 일반 텍스트 특약 어쩌구")["riders"] == []


def test_rider_tokens_filter_out_legal_fragments():
    from core import guide
    # 담보명 섹션에 법조문 조각이 섞여도 담보 키워드 없는 건 제외
    text = ("[메인 담보명]\n특약명\n"
            "항공기 지연 결항 보상(지수형)(국내 출국) 특약\n"
            "이 특약\n① 이 특약\n제1조(특약\n된 경우에는 이 특약\n"
            "여행중 휴대품 손해(분실제외) 특약\n")
    riders = guide.parse_guide(text)["riders"]
    assert "항공기 지연 결항 보상(지수형)(국내 출국) 특약" in riders
    assert "여행중 휴대품 손해(분실제외) 특약" in riders
    assert "이 특약" not in riders
    assert not any(r in ("① 이 특약", "제1조(특약", "된 경우에는 이 특약") for r in riders)


def test_default_riders_constant():
    from core import guide
    assert len(guide.DEFAULT_RIDERS) == 6
    assert "여행중 여권분실 재발급비용 특약" in guide.DEFAULT_RIDERS
